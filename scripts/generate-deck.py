"""
WireTrust Presentation Deck — Premium Design
Run: python scripts/generate-deck.py
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Design System ───────────────────────────────────────────────────────
BG_DARK = RGBColor(0x09, 0x09, 0x0E)
BG_CARD = RGBColor(0x13, 0x13, 0x1A)
BG_CARD_ACCENT = RGBColor(0x16, 0x13, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xE4, 0xE4, 0xE7)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
MUTED = RGBColor(0x71, 0x71, 0x7A)
DARK = RGBColor(0x3F, 0x3F, 0x46)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_GLOW = RGBColor(0x2E, 0x1A, 0x5E)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_DIM = RGBColor(0x0D, 0x2A, 0x1F)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x38, 0x7C, 0xFD)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

def add_shape(slide, left, top, w, h, color=BG_CARD, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_accent_line(slide, left, top, w, color=PURPLE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False

def text(slide, left, top, w, h, txt, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.alignment = align
    return tf

def para(tf, txt, size=16, color=GRAY, bold=False, space=6):
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.space_before = Pt(space)
    return p

def add_glow_circle(slide, x, y, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False

def add_num_badge(slide, x, y, num, color=PURPLE):
    s = add_shape(slide, x, y, 0.5, 0.5, color)
    text(slide, x, y + 0.05, 0.5, 0.4, str(num), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: HERO
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

# Glow effects
add_glow_circle(s, 7, -1, 6, RGBColor(0x1A, 0x0A, 0x3E))
add_glow_circle(s, -1, 3, 5, RGBColor(0x0A, 0x2E, 0x1A))

# Badge
add_shape(s, 1, 1.2, 3.5, 0.45, RGBColor(0x1A, 0x1A, 0x25))
text(s, 1.15, 1.22, 3.2, 0.4, "Live on WireFluid  |  Chain 92533", size=11, color=PURPLE_LIGHT)

text(s, 1, 2.1, 10, 1.2, "AI-Powered Cricket", size=52, color=WHITE, bold=True)
text(s, 1, 3.1, 10, 1, "Fan Economy, On-Chain.", size=52, color=PURPLE_LIGHT, bold=True)

tf = text(s, 1, 4.5, 7, 1.5, "Predict matches, build dream squads, deploy ML agents, earn NFT rewards.", size=20, color=GRAY)
para(tf, "Every action is policy-enforced. Every achievement is permanent.", size=18, color=MUTED)
para(tf, "Free to play. Halal compliant. Powered by WireFluid.", size=18, color=MUTED)

# Stats bar
for i, (val, label) in enumerate([("9", "Smart Contracts"), ("203", "Passing Tests"), ("200", "RF Trees"), ("4", "ML Algorithms")]):
    x = 1 + i * 2.2
    add_shape(s, x, 6.3, 1.9, 0.8, BG_CARD)
    text(s, x, 6.35, 1.9, 0.4, val, size=22, color=PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 6.75, 1.9, 0.3, label, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, RED)
text(s, 1, 1, 10, 0.6, "THE PROBLEM", size=13, color=RED, bold=True)
text(s, 1, 1.5, 10, 0.8, "Cricket fans deserve better", size=38, color=WHITE, bold=True)

problems = [
    ("No Ownership", "Fan predictions, achievements, and engagement history are locked in centralized databases. Platforms own your data.", RED),
    ("Gambling Barrier", "Most engagement platforms rely on betting. This excludes halal-conscious fans across Pakistan and the Muslim world.", AMBER),
    ("Zero Intelligence", "Franchise teams lack accessible ML analytics. Decisions made on gut feeling, not data-driven intelligence.", BLUE),
    ("Trust Deficit", "Scores, streaks, and rewards have no verifiable proof. Rules change without notice. No audit trail.", PURPLE),
]

for i, (title, desc, color) in enumerate(problems):
    x = 0.7 + i * 3.1
    add_shape(s, x, 2.8, 2.8, 4)
    add_accent_line(s, x, 2.8, 2.8, color)
    text(s, x + 0.3, 3.2, 2.2, 0.5, title, size=18, color=color, bold=True)
    text(s, x + 0.3, 3.9, 2.2, 2.5, desc, size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, GREEN)
text(s, 1, 1, 10, 0.6, "THE SOLUTION", size=13, color=GREEN, bold=True)
text(s, 1, 1.5, 10, 0.8, "Four ways to engage", size=38, color=WHITE, bold=True)

solutions = [
    ("Match Predictions", "FREE", "Call match winners, top scorers, total runs. Earn points and streak bonuses. Early bird bonus for quick calls.", GREEN, "Points only. Zero staking."),
    ("Squad Challenge", "FREE", "Pick 11 players within 100 credits. Captain 2x points, Vice-Captain 1.5x. Sponsor-funded prize pools.", BLUE, "Fans never pay to play."),
    ("NFT Rewards", "EARNED", "Match tickets, player collectibles, VIP experiences, soulbound achievement badges. Anti-scalping 110% resale cap.", AMBER, "On-chain, verifiable, tradeable."),
    ("AI Agents", "UNLOCK", "Deploy ML-powered agents after reaching milestones. Random Forest predictions, form analysis, squad optimization.", PURPLE, "Achievement-gated progression."),
]

for i, (title, badge, desc, color, footer) in enumerate(solutions):
    x = 0.7 + i * 3.1
    add_shape(s, x, 2.8, 2.8, 4.2)
    add_accent_line(s, x, 2.8, 2.8, color)
    add_shape(s, x + 0.3, 3.15, 0.8, 0.3, color)
    text(s, x + 0.35, 3.15, 0.7, 0.3, badge, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 3.65, 2.2, 0.5, title, size=17, color=WHITE, bold=True)
    text(s, x + 0.3, 4.25, 2.2, 2, desc, size=12, color=GRAY)
    text(s, x + 0.3, 6.3, 2.2, 0.4, footer, size=10, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: WIREFLUID UTILIZATION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, PURPLE)
text(s, 1, 1, 10, 0.6, "WIREFLUID NETWORK", size=13, color=PURPLE_LIGHT, bold=True)
text(s, 1, 1.5, 10, 0.8, "How WireTrust uses the blockchain", size=38, color=WHITE, bold=True)

uses = [
    ("9 Smart Contracts", "Core pipeline: FranchiseRegistry > AgentRegistry > PolicyEngine > ExecutionGateway > ReputationStore. Plus PredictionModule, FantasyModule, WireTrustNFT, MatchOracle."),
    ("8-Check Policy Engine", "Every agent action validated: ownership, activation, forbidden targets, nonce replay, spending limit, daily cap, frequency, max positions. All on-chain."),
    ("On-Chain Reputation", "Score 0-100 updated per execution. Success +2, failure -5, violation -10. SAFE/MEDIUM/RISKY badges. Permanent, immutable audit trail."),
    ("ERC-721 NFT System", "5 categories with category-specific rules. Soulbound badges. Buyer-initiated marketplace (listForSale + buyToken). 2.5% protocol fee."),
    ("Transparent Settlement", "MatchOracle with multi-oracle pattern. Paginated prediction resolution (no gas limit issues). Pull-pattern prize claims."),
    ("Emergency Controls", "Pausable on 3 critical contracts (Gateway, Fantasy, Prediction). Emergency ETH withdraw from stuck contests. FallbackProvider for RPC."),
]

for i, (title, desc) in enumerate(uses):
    col = i % 2; row = i // 2
    x = 0.7 + col * 6.2; y = 2.6 + row * 1.55
    add_shape(s, x, y, 5.8, 1.35)
    add_num_badge(s, x + 0.2, y + 0.2, i + 1)
    text(s, x + 0.9, y + 0.2, 4.5, 0.4, title, size=15, color=PURPLE_LIGHT, bold=True)
    text(s, x + 0.9, y + 0.65, 4.7, 0.7, desc, size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: SMART CONTRACT ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, CYAN)
text(s, 1, 1, 10, 0.6, "ARCHITECTURE", size=13, color=CYAN, bold=True)
text(s, 1, 1.5, 10, 0.8, "9 contracts, 3,400+ lines Solidity", size=38, color=WHITE, bold=True)

# Left: Core
text(s, 1, 2.6, 5, 0.4, "Core Pipeline (5 contracts)", size=16, color=PURPLE_LIGHT, bold=True)
core = [
    ("FranchiseRegistry", "Multi-tenant franchise onboarding"),
    ("AgentRegistry", "Bot identity, ownership, franchise-scoped"),
    ("ReputationStore", "Behavioral scoring 0-100, risk badges"),
    ("PolicyEngine", "8-check execution constraints"),
    ("ExecutionGateway", "Central router with Pausable + fee collection"),
]
for i, (name, desc) in enumerate(core):
    y = 3.1 + i * 0.8
    add_shape(s, 1, y, 5.5, 0.65)
    text(s, 1.3, y + 0.05, 2.5, 0.35, name, size=13, color=WHITE, bold=True)
    text(s, 1.3, y + 0.35, 4.8, 0.3, desc, size=10, color=MUTED)

# Right: Modules
text(s, 7, 2.6, 5, 0.4, "Modules + Oracle (4 contracts)", size=16, color=GREEN, bold=True)
mods = [
    ("PredictionModule", "Free predictions, streaks, paginated batch resolution"),
    ("FantasyModule", "Sponsor-funded contests, 200 participant cap, pull-pattern prizes"),
    ("WireTrustNFT", "ERC-721, soulbound badges, buyer marketplace, 110% resale cap"),
    ("MatchOracle", "Multi-oracle settlement, player stats, abandon handling"),
]
for i, (name, desc) in enumerate(mods):
    y = 3.1 + i * 0.8
    add_shape(s, 7, y, 5.5, 0.65)
    text(s, 7.3, y + 0.05, 2.5, 0.35, name, size=13, color=WHITE, bold=True)
    text(s, 7.3, y + 0.35, 4.8, 0.3, desc, size=10, color=MUTED)

# Security bar
add_shape(s, 1, 6.4, 11.5, 0.7, GREEN_DIM)
text(s, 1.3, 6.45, 11, 0.6, "Security: ReentrancyGuard + Pausable + Nonce Replay + Fee-After-Success + Soulbound + Anti-Scalping + Participant Cap + Emergency Withdraw", size=11, color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: ML ENGINE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_glow_circle(s, 9, 0, 5, PURPLE_GLOW)
add_accent_line(s, 1, 0.8, 1.5, PURPLE)
text(s, 1, 1, 10, 0.6, "MACHINE LEARNING", size=13, color=PURPLE_LIGHT, bold=True)
text(s, 1, 1.5, 10, 0.8, "Real ML, not heuristics", size=38, color=WHITE, bold=True)

# 6-Factor grid
factors = [
    ("ELO Ratings", "25%", "Self-learning power rankings"),
    ("EWMA Form", "20%", "Exponential weighted averages"),
    ("Head-to-Head", "15%", "Historical matchup dominance"),
    ("Momentum", "15%", "Win streaks, recency-weighted"),
    ("Venue Stats", "15%", "Ground-specific win rates"),
    ("Role-Weighted", "10%", "Batting vs bowling split"),
]
for i, (name, weight, desc) in enumerate(factors):
    col = i % 3; row = i // 3
    x = 0.7 + col * 2.2; y = 2.7 + row * 1.3
    add_shape(s, x, y, 2, 1.05)
    text(s, x + 0.2, y + 0.1, 1.1, 0.3, name, size=11, color=WHITE, bold=True)
    text(s, x + 1.3, y + 0.1, 0.5, 0.3, weight, size=11, color=PURPLE_LIGHT, bold=True)
    text(s, x + 0.2, y + 0.5, 1.6, 0.4, desc, size=9, color=MUTED)

# ML Algorithms
ml = [
    ("Random Forest Classifier", "200 trees", "80/20 chronological train/test split. 6 features per match. Reports held-out TEST accuracy.", PURPLE),
    ("Weighted Linear Regression", "Player forecast", "Exponential decay weighting. Confidence intervals from R-squared + sample size.", GREEN),
    ("Constrained Knapsack", "Squad optimizer", "Maximize predicted FP within 100 credits. Role constraints enforced. Auto captain/VC.", BLUE),
    ("Z-Score Detection", "Anomaly alerts", "BREAKOUT (z>2.0) and COLLAPSE (z<-2.0) with percentile ranking via normal CDF.", AMBER),
]
for i, (name, sub, desc, color) in enumerate(ml):
    x = 7.2; y = 2.7 + i * 1.1
    add_shape(s, x, y, 5.3, 0.9)
    add_accent_line(s, x, y, 0.15, color)
    text(s, x + 0.3, y + 0.08, 3, 0.3, name, size=12, color=color, bold=True)
    text(s, x + 3.5, y + 0.08, 1.5, 0.3, sub, size=10, color=MUTED)
    text(s, x + 0.3, y + 0.42, 4.8, 0.4, desc, size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: FAN JOURNEY
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, GREEN)
text(s, 1, 1, 10, 0.6, "USER JOURNEY", size=13, color=GREEN, bold=True)
text(s, 1, 1.5, 10, 0.8, "From visitor to AI agent owner", size=38, color=WHITE, bold=True)

steps = [
    ("01", "Connect", "Google Sign-In\nor MetaMask.\nFree WIRE credits\non testnet.", GREEN),
    ("02", "Predict", "Call match winners,\ntop scorers, totals.\nEarn points.\nBuild streaks.", BLUE),
    ("03", "Build Squad", "Pick 11 players.\n100 credit budget.\nCaptain 2x.\nWin sponsor prizes.", PURPLE),
    ("04", "Earn NFTs", "Complete challenges.\nTickets, badges,\ncollectibles, VIP\nexperiences.", AMBER),
    ("05", "Unlock AI", "5 predictions +\n1 squad + 100 pts.\nDeploy your own\nML-powered agent.", PURPLE_LIGHT),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_shape(s, x, 2.8, 2.2, 4.2)
    add_accent_line(s, x, 2.8, 2.2, color)
    text(s, x + 0.25, 3.1, 0.6, 0.5, num, size=28, color=color, bold=True)
    text(s, x + 0.25, 3.65, 1.8, 0.4, title, size=18, color=WHITE, bold=True)
    text(s, x + 0.25, 4.2, 1.8, 2.5, desc, size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: FRANCHISE AI
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_glow_circle(s, 8, -1, 6, PURPLE_GLOW)
add_accent_line(s, 1, 0.8, 1.5, PURPLE)
text(s, 1, 1, 10, 0.6, "FRANCHISE INTELLIGENCE", size=13, color=PURPLE_LIGHT, bold=True)
text(s, 1, 1.5, 10, 0.8, "ML agents for team management", size=38, color=WHITE, bold=True)

agents = [
    ("Opposition Scout", "Every 30 min", "Monitors all rival teams for form changes, momentum shifts, and player breakouts. Classifies threats as HIGH/MEDIUM/LOW. Alerts on hot streaks and collapses.", "eye"),
    ("Squad Form Monitor", "Every 15 min", "Tracks squad players using weighted linear regression. Z-score anomaly detection identifies BREAKOUT and COLLAPSE events. Confidence-adjusted predictions.", "activity"),
    ("Match Preparation", "Every 60 min", "Auto-generates scouting reports using Random Forest. 6-feature win probability with confidence intervals. AI tactical notes and recommended XI.", "brain"),
]

for i, (name, interval, desc, icon) in enumerate(agents):
    y = 2.6 + i * 1.5
    add_shape(s, 0.7, y, 11.8, 1.3)
    add_accent_line(s, 0.7, y, 0.15, PURPLE)
    text(s, 1.1, y + 0.15, 4, 0.4, name, size=18, color=WHITE, bold=True)
    add_shape(s, 9.5, y + 0.15, 1.5, 0.3, PURPLE_GLOW)
    text(s, 9.5, y + 0.15, 1.5, 0.3, interval, size=10, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    text(s, 1.1, y + 0.6, 10.5, 0.6, desc, size=12, color=GRAY)

add_shape(s, 0.7, 7, 11.8, 0.3, GREEN_DIM)
text(s, 1, 7, 11, 0.3, "Squad Optimizer: Constrained knapsack finds optimal XI using ML-predicted fantasy points. Auto-selects captain and vice-captain.", size=10, color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: SECURITY & TESTING
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, GREEN)
text(s, 1, 1, 10, 0.6, "SECURITY & TESTING", size=13, color=GREEN, bold=True)
text(s, 1, 1.5, 10, 0.8, "Production-grade hardening", size=38, color=WHITE, bold=True)

# Left column
text(s, 1, 2.6, 5, 0.4, "Smart Contract Security", size=16, color=PURPLE_LIGHT, bold=True)
sc_items = ["ReentrancyGuard on all ETH contracts", "Pausable emergency stop (3 contracts)", "Nonce replay protection", "Fee collected after success only", "Participant cap (200) prevents gas attacks", "Paginated resolution (no unbounded loops)", "Soulbound badges + anti-scalping 110%", "Emergency withdraw for stuck ETH"]
for i, item in enumerate(sc_items):
    text(s, 1.3, 3.1 + i * 0.45, 5, 0.35, f"  {item}", size=11, color=GRAY)

# Right column
text(s, 7, 2.6, 5, 0.4, "Backend Security", size=16, color=GREEN, bold=True)
be_items = ["3-tier rate limiting (strict/moderate/relaxed)", "CSRF via X-Requested-With header", "CORS restricted to configured origins", "Franchise-scoped agent authorization", "Input validation on all mutations", "Global error handler + 404 route", "RPC FallbackProvider", "Leaderboard caching (60s TTL)"]
for i, item in enumerate(be_items):
    text(s, 7.3, 3.1 + i * 0.45, 5, 0.35, f"  {item}", size=11, color=GRAY)

# Test bar
add_shape(s, 0.7, 6.7, 11.8, 0.5, GREEN_DIM)
text(s, 1.1, 6.75, 2, 0.4, "203 TESTS", size=20, color=GREEN, bold=True)
text(s, 3.5, 6.8, 9, 0.35, "Core (57) + Modules (77) + Oracle (13) + Integration (12) + NewFeatures (30) + Gas benchmarks", size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: DEPLOYED CONTRACTS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, CYAN)
text(s, 1, 1, 10, 0.6, "ON-CHAIN", size=13, color=CYAN, bold=True)
text(s, 1, 1.5, 10, 0.8, "Deployed on WireFluid (Chain 92533)", size=38, color=WHITE, bold=True)

deployed = [
    ("FranchiseRegistry", "0x0352a16cd7B6b06707A363452d7f1937840E90D2", PURPLE_LIGHT),
    ("AgentRegistry", "0x603F8Db0b75d7dc699335a8C45412EC2eE49A60c", PURPLE_LIGHT),
    ("ReputationStore", "0x096A05F4F2d81617Ac93589958f9736B3DF2c915", PURPLE_LIGHT),
    ("PolicyEngine", "0x0802bBe692e21ed0f7C00d7bD282ec5fAB2E582C", PURPLE_LIGHT),
    ("ExecutionGateway", "0xC148f0AE3d83089217d33AA2b0b00e1F2b9e889e", PURPLE_LIGHT),
    ("MatchOracle", "0x3c1D9725eD4D92484Ac45af00a0b836C95a76E86", GREEN),
    ("PredictionModule", "0xe175236fe8978FdEDB7d563da3498c12531b241c", GREEN),
    ("FantasyModule", "0x726fEb1Bdf8E4CE4ABcd273C5c1696Ed24f31d30", GREEN),
    ("WireTrustNFT", "0x9127EFFB479ae271601d18Bfb7CF6Af491244e1b", AMBER),
]

for i, (name, addr, color) in enumerate(deployed):
    y = 2.5 + i * 0.52
    add_shape(s, 0.7, y, 11.8, 0.42)
    text(s, 1.1, y + 0.05, 3, 0.3, name, size=13, color=color, bold=True)
    text(s, 4.2, y + 0.05, 8, 0.3, addr, size=12, color=MUTED)

add_shape(s, 0.7, 7.3, 11.8, 0.42, BG_CARD)
text(s, 1.1, 7.35, 3, 0.3, "Deployer / Treasury", size=13, color=WHITE, bold=True)
text(s, 4.2, 7.35, 8, 0.3, "0x22EfFAe93649A93F7c6e01aBB6Ce2496BB2D4105", size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: REAL WORLD IMPACT
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, AMBER)
text(s, 1, 1, 10, 0.6, "IMPACT", size=13, color=AMBER, bold=True)
text(s, 1, 1.5, 10, 0.8, "Why this matters", size=38, color=WHITE, bold=True)

impacts = [
    ("Halal Fan Engagement", "No betting, no staking, no gambling. Points-only predictions with sponsor-funded prizes. Designed for 2 billion Muslims worldwide. Zero entry barriers.", GREEN),
    ("Cricket Fan Ownership", "Every prediction, squad entry, streak, and NFT reward is permanently recorded on WireFluid. Fans own their history. Platforms cannot erase achievements.", PURPLE),
    ("Franchise Intelligence", "PSL teams get ML scouting, player form forecasting, anomaly detection, squad optimization. Data-driven decisions instead of gut feeling.", BLUE),
    ("Free Economy", "Sign in with Google. No crypto knowledge needed. Free testnet credits. Web3Auth abstracts the blockchain. Sponsor-funded prizes, not fan fees.", AMBER),
]

for i, (title, desc, color) in enumerate(impacts):
    col = i % 2; row = i // 2
    x = 0.7 + col * 6.2; y = 2.6 + row * 2.3
    add_shape(s, x, y, 5.8, 2)
    add_accent_line(s, x, y, 5.8, color)
    text(s, x + 0.3, y + 0.35, 5.2, 0.5, title, size=18, color=color, bold=True)
    text(s, x + 0.3, y + 0.9, 5.2, 1.1, desc, size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: TECH STACK
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_accent_line(s, 1, 0.8, 1.5, BLUE)
text(s, 1, 1, 10, 0.6, "TECHNOLOGY", size=13, color=BLUE, bold=True)
text(s, 1, 1.5, 10, 0.8, "Full-stack architecture", size=38, color=WHITE, bold=True)

stack = [
    ("Blockchain", "WireFluid EVM (Chain 92533), Solidity ^0.8.27, OpenZeppelin v5, Hardhat", PURPLE_LIGHT),
    ("Contracts", "9 deployed: 5 core + 3 modules + 1 oracle. Pausable + ReentrancyGuard", PURPLE_LIGHT),
    ("Backend", "Node.js, Express, ethers.js v6, PostgreSQL, WebSocket, PM2", GREEN),
    ("Frontend", "React 18, Vite, Tailwind CSS, Redux Toolkit, Web3Auth + MetaMask", BLUE),
    ("ML Engine", "ml-random-forest (200 trees), weighted regression, knapsack, z-score", AMBER),
    ("Security", "Rate limiting, CSRF, CORS, RPC failover, Pausable, input validation", RED),
    ("Testing", "203 tests across 5 suites. Hardhat + Chai + ethers. Gas reporter.", GREEN),
    ("Data", "157 PSL players, 134 matches, 6 franchise teams, real ESPNcricinfo data", CYAN),
]

for i, (label, desc, color) in enumerate(stack):
    y = 2.5 + i * 0.6
    add_shape(s, 0.7, y, 11.8, 0.48)
    text(s, 1.1, y + 0.08, 2, 0.35, label, size=13, color=color, bold=True)
    text(s, 3.3, y + 0.08, 9, 0.35, desc, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: CLOSING
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

add_glow_circle(s, 4, 0.5, 5, PURPLE_GLOW)
add_glow_circle(s, 6, 3, 4, RGBColor(0x0A, 0x2E, 0x1A))

text(s, 1, 2.2, 11.3, 1, "WireTrust", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, 1, 3.3, 11.3, 0.6, "AI-Powered Cricket Fan Economy on WireFluid", size=24, color=PURPLE_LIGHT, align=PP_ALIGN.CENTER)

add_shape(s, 4, 4.3, 5.3, 0.45, BG_CARD)
text(s, 4, 4.33, 5.3, 0.4, "https://wiretrust.metagenesis.work", size=15, color=GREEN, align=PP_ALIGN.CENTER)

add_shape(s, 4, 4.9, 5.3, 0.45, BG_CARD)
text(s, 4, 4.93, 5.3, 0.4, "github.com/hamzzaaamalik/wiretrust", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom stats
for i, (val, label) in enumerate([("9", "Contracts"), ("203", "Tests"), ("200", "RF Trees"), ("30", "Pages"), ("80+", "Endpoints")]):
    x = 2 + i * 1.9
    text(s, x, 5.9, 1.5, 0.4, val, size=22, color=PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 6.3, 1.5, 0.3, label, size=10, color=MUTED, align=PP_ALIGN.CENTER)

text(s, 1, 7, 11.3, 0.3, "Free to Play  |  Halal Compliant  |  On-Chain Verified  |  ML-Powered", size=12, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
output = "WireTrust_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
